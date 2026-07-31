"""
build_ppt.py — Generate viva presentation PPTX for the NTCC project.
Uses python-pptx (since Node is unavailable on this Windows box).

Style: Vintage & Academic palette (#4) - classic, scholarly.
Output: NTCC_Viva_Presentation.pptx in the project root.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ── Theme (Vintage & Academic palette) ──────────────────────────────────────
DEEP_BLUE    = RGBColor(0x00, 0x30, 0x49)   # primary text / titles
LIGHT_BLUE   = RGBColor(0x66, 0x9b, 0xbc)   # secondary text / accents
BRIGHT_RED   = RGBColor(0xc1, 0x12, 0x1f)   # highlights, key numbers
CREAM        = RGBColor(0xfd, 0xf0, 0xd5)   # card backgrounds
DEEP_RED     = RGBColor(0x78, 0x00, 0x00)   # secondary accent
WHITE        = RGBColor(0xff, 0xff, 0xff)
DARK_TEXT    = RGBColor(0x1a, 0x1a, 0x1a)
GREY_TEXT    = RGBColor(0x55, 0x55, 0x55)
LIGHT_GREY   = RGBColor(0xee, 0xee, 0xee)
MID_GREY     = RGBColor(0xaa, 0xaa, 0xaa)

# Layout: standard widescreen 16:9 (13.333" x 7.5")
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Fonts
FONT_TITLE = "Calibri"
FONT_BODY  = "Calibri"
FONT_MONO  = "Consolas"

OUT_PATH = Path(__file__).resolve().parent / "NTCC_Viva_Presentation.pptx"

# ── Helpers ────────────────────────────────────────────────────────────────

def _set_slide_bg(slide, rgb):
    """Solid-fill the slide background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _add_rect(slide, x, y, w, h, fill, line=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Pt(0.75)
    rect.shadow.inherit = False
    return rect


def _add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK_TEXT,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_BODY, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def _add_bullets(slide, x, y, w, h, items, *, size=16, color=DARK_TEXT,
                 bullet_color=DEEP_BLUE, line_spacing=1.25, bold_first=False):
    """items: list of (text, sub_items) tuples or just strings."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, subs = item
        else:
            text, subs = item, []
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(3)
        # Bullet marker
        bullet_run = p.add_run()
        bullet_run.text = "▸  "
        bullet_run.font.name = FONT_BODY
        bullet_run.font.size = Pt(size)
        bullet_run.font.bold = True
        bullet_run.font.color.rgb = bullet_color
        # Text
        run = p.add_run()
        run.text = text
        run.font.name = FONT_BODY
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold_first and i == 0
        # Sub-items
        for sub in subs:
            sp = tf.add_paragraph()
            sp.alignment = PP_ALIGN.LEFT
            sp.line_spacing = line_spacing
            sp.space_after = Pt(2)
            sp.level = 1
            sb = sp.add_run()
            sb.text = "•   "
            sb.font.name = FONT_BODY
            sb.font.size = Pt(size - 2)
            sb.font.color.rgb = MID_GREY
            sr = sp.add_run()
            sr.text = sub
            sr.font.name = FONT_BODY
            sr.font.size = Pt(size - 2)
            sr.font.color.rgb = GREY_TEXT
    return tb


def _add_page_number(slide, n, total):
    tb = slide.shapes.add_textbox(Inches(12.4), Inches(7.1), Inches(0.85), Inches(0.3))
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{n} / {total}"
    r.font.name = FONT_BODY
    r.font.size = Pt(10)
    r.font.color.rgb = MID_GREY


def _add_header_bar(slide, title, subtitle=None):
    # Top blue bar
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.55), DEEP_BLUE)
    # Title text
    _add_text(slide, Inches(0.4), Inches(0.07), Inches(11), Inches(0.5),
              title, size=22, bold=True, color=WHITE, font=FONT_TITLE)
    if subtitle:
        _add_text(slide, Inches(0.4), Inches(0.6), Inches(12.5), Inches(0.35),
                  subtitle, size=12, color=DEEP_BLUE, italic=True, font=FONT_BODY)
    # Red accent strip below header
    _add_rect(slide, Inches(0), Inches(0.55), SLIDE_W, Inches(0.06), BRIGHT_RED)


def _add_footer(slide):
    _add_text(slide, Inches(0.4), Inches(7.15), Inches(8), Inches(0.3),
              "Semantic Image Search Engine using CLIP | NTCC Viva | Dhruv Singhal",
              size=9, color=MID_GREY, italic=True)


def _add_section_divider(slide, x, y, w, color=BRIGHT_RED):
    _add_rect(slide, x, y, w, Inches(0.04), color)


# ── Slide builders ─────────────────────────────────────────────────────────

def _cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(s, DEEP_BLUE)
    # Big red accent stripe
    _add_rect(s, Inches(0), Inches(2.4), SLIDE_W, Inches(0.1), BRIGHT_RED)
    # Top deep red band
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.5), DEEP_RED)
    _add_text(s, Inches(0.5), Inches(0.1), Inches(10), Inches(0.3),
              "AMITY UNIVERSITY, NOIDA  •  BTECH CSE  •  NTCC VIVA 2026",
              size=11, bold=True, color=WHITE, font=FONT_TITLE)
    # Title
    _add_text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(1.0),
              "Semantic Image Search Engine",
              size=44, bold=True, color=WHITE, font=FONT_TITLE)
    _add_text(s, Inches(0.7), Inches(1.65), Inches(12), Inches(0.7),
              "using CLIP for Cross-Modal Retrieval and Vision-Language Alignment",
              size=22, color=LIGHT_BLUE, italic=True, font=FONT_TITLE)
    # Author block
    _add_text(s, Inches(0.7), Inches(2.9), Inches(8), Inches(0.5),
              "Presented by:", size=14, color=LIGHT_BLUE, font=FONT_BODY)
    _add_text(s, Inches(0.7), Inches(3.25), Inches(8), Inches(0.6),
              "DHRUV SINGHAL", size=28, bold=True, color=WHITE, font=FONT_TITLE)
    _add_text(s, Inches(0.7), Inches(3.8), Inches(8), Inches(0.4),
              "Enrollment No: A2305224486  |  BTECH CSE 2024–2028",
              size=14, color=LIGHT_BLUE, font=FONT_BODY)
    # Supervisor block
    _add_text(s, Inches(0.7), Inches(4.6), Inches(8), Inches(0.4),
              "Under the guidance of:", size=14, color=LIGHT_BLUE, font=FONT_BODY)
    _add_text(s, Inches(0.7), Inches(4.95), Inches(10), Inches(0.5),
              "Dr. Rakesh Chandra Joshi", size=22, bold=True, color=WHITE, font=FONT_TITLE)
    _add_text(s, Inches(0.7), Inches(5.45), Inches(10), Inches(0.4),
              "Assistant Professor, Amity Centre for Artificial Intelligence",
              size=13, color=LIGHT_BLUE, italic=True, font=FONT_BODY)
    # Right card with key info
    _add_rect(s, Inches(9.0), Inches(2.9), Inches(3.8), Inches(3.0), WHITE)
    _add_rect(s, Inches(9.0), Inches(2.9), Inches(3.8), Inches(0.45), BRIGHT_RED)
    _add_text(s, Inches(9.15), Inches(2.95), Inches(3.6), Inches(0.35),
              "PROJECT AT A GLANCE", size=12, bold=True, color=WHITE, font=FONT_TITLE)
    info_items = [
        ("Datasets", "Flickr8k · CC3M · RF/Telemetry"),
        ("Models",   "CLIP · BLIP · ALIGN"),
        ("Search modes", "T2I · I2T · I2I · Semantic vs Lexical"),
        ("Deployment", "ONNX + INT8, CPU <100 ms"),
        ("Live app", "image-search-app.streamlit.app"),
    ]
    for i, (k, v) in enumerate(info_items):
        y = Inches(3.5 + i * 0.45)
        _add_text(s, Inches(9.15), y, Inches(1.5), Inches(0.35), k,
                  size=11, bold=True, color=DEEP_BLUE, font=FONT_TITLE)
        _add_text(s, Inches(9.15), y + Inches(0.22), Inches(3.5), Inches(0.3), v,
                  size=10, color=DARK_TEXT, font=FONT_BODY)
    # Bottom band
    _add_rect(s, Inches(0), Inches(7.0), SLIDE_W, Inches(0.5), DEEP_RED)
    _add_text(s, Inches(0.5), Inches(7.1), Inches(12), Inches(0.3),
              "In Partial Fulfilment of the Requirement for the Degree of Bachelor of Technology in Computer Science and Engineering",
              size=11, color=WHITE, italic=True, font=FONT_BODY)


def _outline(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(s, WHITE)
    _add_header_bar(s, "Presentation Outline",
                    "15-slide walkthrough of motivation, methodology, results, and live demo")
    items = [
        ("01", "Motivation & Problem Statement",
         "Vocabulary mismatch → CLIP solution → why domain adaptation matters"),
        ("02", "Research Objectives & Contributions",
         "5 goals + 4 key contributions of this work"),
        ("03", "Related Work",
         "CLIP, ALIGN, BLIP-2, CoOp, MaPLe, MobileCLIP, SigLIP, MARVEL"),
        ("04", "Datasets",
         "Flickr8k (8K RF/telemetry) + CC3M (1K) + OOD-RF"),
        ("05", "Methodology",
         "Dual-tower · 2-layer MLP head · warm-start · InfoNCE loss"),
        ("06", "PEFT & Domain Adaptation",
         "Freeze 99% backbone, train only projection head (0.06–0.4% params)"),
        ("07", "Dense Retrieval & ONNX/INT8",
         "FAISS IndexFlatIP + ONNX export + dynamic INT8 quantisation"),
        ("08", "Results",
         "Precision@K, latency, cross-model benchmark, dense vs sparse, ablations"),
        ("09", "Live Demo",
         "5 tabs of the Streamlit app · spectrogram + text-refinement fallback"),
        ("10", "WPR Bugs & Fixes",
         "5 critical deploy failures and their resolutions"),
        ("11", "Conclusion & Future Scope",
         "6 directions: MobileCLIP, SigLIP, video, continual, hybrid, active"),
    ]
    # 2 columns
    for i, (num, title, desc) in enumerate(items):
        col = i // 6
        row = i % 6
        x = Inches(0.5 + col * 6.3)
        y = Inches(1.0 + row * 1.05)
        # Number circle
        _add_rect(s, x, y, Inches(0.65), Inches(0.65), DEEP_BLUE)
        _add_text(s, x, y, Inches(0.65), Inches(0.65), num,
                  size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                  anchor=MSO_ANCHOR.MIDDLE, font=FONT_TITLE)
        # Title
        _add_text(s, x + Inches(0.8), y + Inches(0.02), Inches(5.3), Inches(0.4),
                  title, size=14, bold=True, color=DEEP_BLUE, font=FONT_TITLE)
        # Description
        _add_text(s, x + Inches(0.8), y + Inches(0.42), Inches(5.3), Inches(0.55),
                  desc, size=10, color=GREY_TEXT, italic=True, font=FONT_BODY)
    _add_footer(s)
    _add_page_number(s, n, total)


def _motivation(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(s, WHITE)
    _add_header_bar(s, "Motivation & Problem Statement",
                    "Why vision-language retrieval needs domain adaptation, not just bigger models")
    # Left column: the problem
    _add_rect(s, Inches(0.4), Inches(1.1), Inches(6.0), Inches(5.7), CREAM)
    _add_rect(s, Inches(0.4), Inches(1.1), Inches(6.0), Inches(0.5), DEEP_RED)
    _add_text(s, Inches(0.6), Inches(1.18), Inches(5.6), Inches(0.4),
              "The Vocabulary Mismatch Problem", size=15, bold=True, color=WHITE, font=FONT_TITLE)
    _add_bullets(s, Inches(0.65), Inches(1.75), Inches(5.5), Inches(4.8), [
        "Traditional search uses captions + BM25 / TF-IDF",
        ("Fails on semantically equivalent but lexically different pairs:", [
            "Query: \"a vivacious furry companion darting across a sandy waterfront\"",
            "Gold caption: \"dog running on beach\"",
            "→ 0 keywords match, BM25 score = 0.0",
        ]),
        "Pre-trained VLMs (CLIP, ALIGN) solve this — 400M image-text pairs",
        "But: trained on natural photos only",
        "Specialised domains (RF, spectrograms, telemetry) remain OOD",
        "Full fine-tuning of ViT-H/14 needs >40 GB VRAM",
        "Naive fine-tuning causes catastrophic forgetting of zero-shot ability",
    ], size=13)
    # Right column: our framing
    _add_rect(s, Inches(6.7), Inches(1.1), Inches(6.2), Inches(5.7), WHITE, line=LIGHT_BLUE)
    _add_rect(s, Inches(6.7), Inches(1.1), Inches(6.2), Inches(0.5), DEEP_BLUE)
    _add_text(s, Inches(6.9), Inches(1.18), Inches(5.8), Inches(0.4),
              "Our Framing", size=15, bold=True, color=WHITE, font=FONT_TITLE)
    _add_bullets(s, Inches(6.95), Inches(1.75), Inches(5.7), Inches(4.8), [
        ("Core problem", []),
        ("Build a cross-modal retrieval engine that:", [
            "Returns Top-K in <100 ms on CPU-only",
            "Supports T2I, I2I, I2T queries",
            "Preserves the encoder's open-vocabulary ability",
            "Adapts to specialised domains (RF/telemetry)",
        ]),
        ("Constraint", []),
        ("Fine-tune on a curated corpus of 8K image-text pairs only", []),
        ("Solution direction", []),
        ("Parameter-efficient fine-tuning (PEFT):", [
            "Freeze 99% of ViT + text transformer",
            "Train only the cross-modal projection head",
            "~590K trainable params (~0.1% of backbone)",
        ]),
    ], size=13, bold_first=True)
    _add_footer(s)
    _add_page_number(s, n, total)


def _objectives(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(s, WHITE)
    _add_header_bar(s, "Research Objectives & Contributions",
                    "5 research goals, 4 key contributions — all backed by the WPR milestones")
    _add_text(s, Inches(0.4), Inches(1.1), Inches(6.0), Inches(0.4),
              "5 Research Objectives", size=16, bold=True, color=DEEP_BLUE, font=FONT_TITLE)
    obj = [
        ("01", "Parameter-Efficient Domain Adaptation", "adapt frozen CLIP to RF/telemetry via projection-head-only training"),
        ("02", "Cross-Modal Benchmarking", "compare CLIP, BLIP-2, ALIGN on retrieval + captioning"),
        ("03", "Dense vs Sparse Retrieval Validation", "show dense latent beats BM25/TF-IDF on zero-overlap queries"),
        ("04", "Edge Deployment & Quantisation", "ONNX + INT8 → <100 ms CPU latency"),
        ("05", "Alternative Loss Analysis", "InfoNCE vs SigLIP — when to switch"),
    ]
    for i, (n_, title, desc) in enumerate(obj):
        y = Inches(1.55 + i * 0.95)
        _add_rect(s, Inches(0.4), y, Inches(0.7), Inches(0.7), BRIGHT_RED)
        _add_text(s, Inches(0.4), y, Inches(0.7), Inches(0.7), n_,
                  size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                  anchor=MSO_ANCHOR.MIDDLE, font=FONT_TITLE)
        _add_text(s, Inches(1.25), y, Inches(5.0), Inches(0.4),
                  title, size=14, bold=True, color=DEEP_BLUE, font=FONT_TITLE)
        _add_text(s, Inches(1.25), y + Inches(0.38), Inches(5.0), Inches(0.5),
                  desc, size=11, color=GREY_TEXT, italic=True, font=FONT_BODY)
    # Right column: contributions
    _add_text(s, Inches(6.9), Inches(1.1), Inches(6.0), Inches(0.4),
              "4 Key Contributions", size=16, bold=True, color=DEEP_BLUE, font=FONT_TITLE)
    contrib = [
        ("End-to-end system", "Streamlit app with 5 search modes + reactive UI + query history"),
        ("PEFT pipeline", "Fine-tune CLIP/BLIP-2/ALIGN projection heads on RF/telemetry, +13.4 pp P@1"),
        ("Dense-vs-sparse proof", "Zero-keyword-overlap benchmark: CLIP 62.3% vs BM25 18.2% vs TF-IDF 12.5%"),
        ("ONNX + INT8 pipeline", "3.1× speedup, 15.5 MB bundle, <70 ms CPU latency"),
    ]
    for i, (title, desc) in enumerate(contrib):
        y = Inches(1.55 + i * 1.2)
        _add_rect(s, Inches(6.9), y, Inches(6.0), Inches(1.0), CREAM)
        _add_rect(s, Inches(6.9), y, Inches(0.12), Inches(1.0), DEEP_BLUE)
        _add_text(s, Inches(7.15), y + Inches(0.1), Inches(5.5), Inches(0.4),
                  title, size=13, bold=True, color=DEEP_BLUE, font=FONT_TITLE)
        _add_text(s, Inches(7.15), y + Inches(0.45), Inches(5.5), Inches(0.5),
                  desc, size=11, color=DARK_TEXT, font=FONT_BODY)
    _add_footer(s)
    _add_page_number(s, n, total)


def _related_work(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(s, WHITE)
    _add_header_bar(s, "Related Work",
                    "10 SOTA papers, organised by the 4 subfields this thesis integrates")
    themes = [
        ("VLM Pre-training",  ["CLIP (Radford 2021) — 400M pairs, dual-tower InfoNCE",
                                "ALIGN (Jia 2021) — noisy labels, ViT-H/14",
                                "BLIP-2 (Li 2023) — Q-Former + frozen LLM"]),
        ("PEFT & Adapters",   ["CLIP-Adapter (Gao 2021) — bottleneck + residual α",
                                "CoOp (Zhou 2022) — learnable soft prompts",
                                "MaPLe (Khattak 2023) — deep branch-aware prompts"]),
        ("Retrieval & Loss",  ["MARVEL (Zhou 2024) — dense > sparse, +5.7% R@10",
                                "SigLIP (Zhai 2023) — pairwise sigmoid, big-batch friendly",
                                "AudioCLIP (Guzhov 2022) — extends to non-standard modalities"]),
        ("Edge Deployment",   ["MobileCLIP (Vasu 2024) — RLHF distillation for edge",
                                "Q-VLM (Wang 2024) — post-training quantisation for VLMs",
                                "Don't Stop Learning (Ding 2023) — continual learning for CLIP"]),
    ]
    cols = 4
    for i, (theme, items) in enumerate(themes):
        x = Inches(0.4 + i * 3.2)
        # Header
        _add_rect(s, x, Inches(1.1), Inches(3.0), Inches(0.55), DEEP_BLUE)
        _add_text(s, x, Inches(1.18), Inches(3.0), Inches(0.4),
                  theme, size=12, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_TITLE)
        # Body card
        _add_rect(s, x, Inches(1.65), Inches(3.0), Inches(5.2), CREAM)
        for j, item in enumerate(items):
            y = Inches(1.85 + j * 1.65)
            # paper title
            _add_text(s, x + Inches(0.15), y, Inches(2.7), Inches(0.4),
                      item.split("—")[0].strip(), size=11, bold=True, color=DEEP_RED, font=FONT_BODY)
            _add_text(s, x + Inches(0.15), y + Inches(0.38), Inches(2.7), Inches(1.2),
                      "— " + item.split("—")[1].strip() if "—" in item else item,
                      size=10, color=DARK_TEXT, font=FONT_BODY)
    _add_text(s, Inches(0.4), Inches(6.95), Inches(12.5), Inches(0.3),
              "Our integration: 4 previously disjoint subfields combined into one working system",
              size=11, bold=True, italic=True, color=DEEP_BLUE, align=PP_ALIGN.CENTER, font=FONT_BODY)
    _add_footer(s)
    _add_page_number(s, n, total)


def _datasets(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(s, WHITE)
    _add_header_bar(s, "Datasets",
                    "3 corpora: Flickr8k (in-domain), CC3M (training), RF/telemetry (specialised)")
    ds = [
        ("Flickr8k — Primary Gallery", "8,000 image-text pairs", DEEP_BLUE, [
            "Standard Flickr30k archive, keyword-filtered for RF/telemetry",
            "Filter keywords: radio, frequency, spectrum, signal, wave, modulation, transmission, antenna, wireless, QAM, PSK, FSK, OFDM, channel, interference",
            "9,200 candidates → uniformly sampled 8,000",
            "Split: 6,400 train / 1,600 held-out validation",
            "Used as the searchable gallery in the live demo (300-img subset for cloud)",
        ]),
        ("CC3M — Conceptual Captions", "1,000 image-text pairs", DEEP_RED, [
            "HuggingFace `pixparse/cc3m-wds` mirror",
            "Subset of 3.3M web-scraped image-caption pairs",
            "Used for initial CPU-based training (small scale)",
            "Demonstrates training loop stability without GPU acceleration",
            "Verifies that the projection-head-only strategy is viable",
        ]),
        ("RF / Telemetry (OOD)",       "Synthetic + real RF captures", BRIGHT_RED, [
            "STFT spectrograms (1024 window, 256 hop, Hanning, viridis)",
            "Generates pseudo-spectrograms from synthetic signals",
            "Used for OOD benchmark in Tab 5 of the demo",
            "Tests graceful-degradation behaviour (\"Relative Semantic Proxy effect\")",
            "Future: extend to birdsong, radar, medical imaging",
        ]),
    ]
    for i, (name, size_lbl, color, items) in enumerate(ds):
        x = Inches(0.4 + i * 4.3)
        # header
        _add_rect(s, x, Inches(1.1), Inches(4.1), Inches(0.9), color)
        _add_text(s, x + Inches(0.15), Inches(1.15), Inches(3.8), Inches(0.4),
                  name, size=13, bold=True, color=WHITE, font=FONT_TITLE)
        _add_text(s, x + Inches(0.15), Inches(1.5), Inches(3.8), Inches(0.45),
                  size_lbl, size=11, color=WHITE, italic=True, font=FONT_BODY)
        # body
        _add_rect(s, x, Inches(2.0), Inches(4.1), Inches(4.85), CREAM)
        _add_bullets(s, x + Inches(0.2), Inches(2.15), Inches(3.7), Inches(4.6), items, size=11, line_spacing=1.15)
    _add_footer(s)
    _add_page_number(s, n, total)


def _methodology(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(s, WHITE)
    _add_header_bar(s, "Methodology — System Architecture",
                    "Dual-tower CLIP with parameter-efficient projection head + dense retrieval")
    # Big architecture diagram (boxes and arrows)
    # Query box
    _add_rect(s, Inches(0.5), Inches(1.2), Inches(2.2), Inches(1.2), DEEP_BLUE)
    _add_text(s, Inches(0.5), Inches(1.25), Inches(2.2), Inches(0.4),
              "QUERY", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_TITLE)
    _add_text(s, Inches(0.5), Inches(1.6), Inches(2.2), Inches(0.7),
              "Text (T2I)\nImage (I2I/T2I)", size=10, color=LIGHT_BLUE, align=PP_ALIGN.CENTER, font=FONT_BODY)
    # Image tower
    _add_rect(s, Inches(3.1), Inches(1.0), Inches(2.3), Inches(1.5), LIGHT_BLUE)
    _add_text(s, Inches(3.1), Inches(1.05), Inches(2.3), Inches(0.4),
              "IMAGE TOWER", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_TITLE)
    _add_text(s, Inches(3.1), Inches(1.4), Inches(2.3), Inches(1.0),
              "ViT-B/32 or ViT-L/14\n(FROZEN)\n→ 512-d feature", size=10, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_BODY)
    # Text tower
    _add_rect(s, Inches(3.1), Inches(2.7), Inches(2.3), Inches(1.5), LIGHT_BLUE)
    _add_text(s, Inches(3.1), Inches(2.75), Inches(2.3), Inches(0.4),
              "TEXT TOWER", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_TITLE)
    _add_text(s, Inches(3.1), Inches(3.1), Inches(2.3), Inches(1.0),
              "Transformer\n(FROZEN)\n→ 512-d feature", size=10, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_BODY)
    # Projection head (TRAINABLE)
    _add_rect(s, Inches(5.9), Inches(1.85), Inches(2.5), Inches(1.5), BRIGHT_RED)
    _add_text(s, Inches(5.9), Inches(1.9), Inches(2.5), Inches(0.4),
              "PROJECTION HEAD", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_TITLE)
    _add_text(s, Inches(5.9), Inches(2.25), Inches(2.5), Inches(1.1),
              "2-layer MLP\n512→384→256\nGELU + dropout 0.1\n(TRAINABLE)", size=10, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_BODY)
    # L2 norm
    _add_rect(s, Inches(8.9), Inches(1.85), Inches(1.6), Inches(1.5), DEEP_BLUE)
    _add_text(s, Inches(8.9), Inches(1.9), Inches(1.6), Inches(0.4),
              "L2 NORM", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_TITLE)
    _add_text(s, Inches(8.9), Inches(2.25), Inches(1.6), Inches(1.1),
              "→ 256-d unit\nvector on\nhypersphere", size=10, color=LIGHT_BLUE, align=PP_ALIGN.CENTER, font=FONT_BODY)
    # FAISS index
    _add_rect(s, Inches(11.0), Inches(1.85), Inches(1.9), Inches(1.5), DEEP_RED)
    _add_text(s, Inches(11.0), Inches(1.9), Inches(1.9), Inches(0.4),
              "FAISS", size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_TITLE)
    _add_text(s, Inches(11.0), Inches(2.25), Inches(1.9), Inches(1.1),
              "IndexFlatIP\nexact NN\nTop-K cosine", size=10, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_BODY)
    # Arrows (text boxes with arrow characters work in PPT)
    for x1, x2, y in [(2.7, 3.1, 1.65), (2.7, 3.1, 3.45), (5.4, 5.9, 2.5), (8.4, 8.9, 2.5), (10.5, 11.0, 2.5)]:
        _add_text(s, Inches(x1), Inches(y - 0.15), Inches(x2 - x1), Inches(0.3),
                  "─" * 5 + "▶", size=12, color=DEEP_BLUE, align=PP_ALIGN.CENTER, font=FONT_MONO)
    # Loss box
    _add_rect(s, Inches(3.1), Inches(4.5), Inches(9.8), Inches(0.8), CREAM, line=DEEP_BLUE)
    _add_text(s, Inches(3.1), Inches(4.55), Inches(9.8), Inches(0.35),
              "Training Loss: Symmetric InfoNCE", size=13, bold=True, color=DEEP_BLUE, align=PP_ALIGN.CENTER, font=FONT_TITLE)
    _add_text(s, Inches(3.1), Inches(4.9), Inches(9.8), Inches(0.35),
              "τ = 0.07  |  AdamW (lr=3e-4, weight_decay=0.01)  |  5 warmup + cosine LR  |  Batch=128  |  FP16 backbone + FP32 Adam states",
              size=11, color=DARK_TEXT, align=PP_ALIGN.CENTER, font=FONT_BODY)
    # Key design choices
    _add_text(s, Inches(0.4), Inches(5.55), Inches(12.5), Inches(0.4),
              "Key Design Choices", size=14, bold=True, color=DEEP_BLUE, font=FONT_TITLE)
    _add_bullets(s, Inches(0.5), Inches(5.95), Inches(12.5), Inches(1.2), [
        ("Warm-start: MLP output initialised to resemble original CLIP embedding projected onto first 256 PCs (under L2)", []),
        ("Symmetric loss = (image-to-text + text-to-image) / 2 → embedding space is consistent for both directions", []),
        ("Cosine similarity on L2-normalised vectors = simple dot product → trivial to FAISS-index and parallelise", []),
    ], size=11, line_spacing=1.1)
    _add_footer(s)
    _add_page_number(s, n, total)


def _peft(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(s, WHITE)
    _add_header_bar(s, "PEFT & Domain Adaptation",
                    "Why we only train 0.06–0.4% of the parameters — and how warm-start saves convergence")
    # Left: parameter breakdown table
    _add_text(s, Inches(0.4), Inches(1.1), Inches(6.0), Inches(0.4),
              "Parameter-Efficient Breakdown", size=14, bold=True, color=DEEP_BLUE, font=FONT_TITLE)
    tbl_data = [
        ["Backbone", "Total params", "Trainable", "% of total"],
        ["ViT-B/32",  "~151 M",      "~590 K",    "0.40%"],
        ["ViT-L/14",  "~430 M",      "~590 K",    "0.14%"],
        ["ViT-H/14",  "~1.0 B",      "~590 K",    "0.06%"],
        ["CLIP-Adapter", "—",        "—",         "≈ 0.10%"],
        ["CoOp",      "—",           "—",         "≈ 0.01%"],
    ]
    table = s.shapes.add_table(len(tbl_data), 4, Inches(0.4), Inches(1.55), Inches(6.0), Inches(2.5)).table
    for r, row in enumerate(tbl_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = FONT_BODY
                    run.font.size = Pt(11)
                    run.font.bold = (r == 0)
                    run.font.color.rgb = WHITE if r == 0 else DARK_TEXT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DEEP_BLUE
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CREAM
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    # Right: warm-start
    _add_text(s, Inches(6.7), Inches(1.1), Inches(6.2), Inches(0.4),
              "Warm-Start Initialisation", size=14, bold=True, color=DEEP_BLUE, font=FONT_TITLE)
    _add_rect(s, Inches(6.7), Inches(1.55), Inches(6.2), Inches(2.5), CREAM)
    _add_bullets(s, Inches(6.85), Inches(1.7), Inches(5.9), Inches(2.3), [
        ("At t=0, the new MLP output is initialised to look like the original CLIP embedding projected onto its first 256 principal components (under L2 norm).", []),
        ("This makes the head an \"identity-like\" mapping within the original subspace.", []),
        ("Training then smoothly reshapes the geometry without destabilising gradients.", []),
        ("Empirically: enables +9.7 pp gain at P@1 from PEFT alone (vs +0 if random init).", []),
    ], size=11, line_spacing=1.2)
    # Bottom: 8K subset
    _add_rect(s, Inches(0.4), Inches(4.4), Inches(12.5), Inches(2.5), DEEP_BLUE)
    _add_text(s, Inches(0.6), Inches(4.5), Inches(12), Inches(0.4),
              "Why This Matters in Practice", size=14, bold=True, color=WHITE, font=FONT_TITLE)
    _add_bullets(s, Inches(0.6), Inches(4.95), Inches(12), Inches(1.9), [
        ("Trains comfortably on Kaggle dual-T4 (2 × 16 GB) — peak VRAM 9.8 GB for L/14, 14.2 GB for H/14", []),
        ("Avoids catastrophic forgetting: the frozen FFN blocks preserve out-of-distribution zero-shot ability", []),
        ("Conforms to CLIP-Adapter, CoOp, MaPLe precedent — rigorously validated in the literature", []),
        ("2-layer MLP > 1-layer linear: +2.5 pp at P@1 (MaPLe thesis confirmed: deeper non-linear projections capture richer cross-modal alignment)", []),
    ], size=12, color=WHITE, bullet_color=BRIGHT_RED)
    _add_footer(s)
    _add_page_number(s, n, total)


def _retrieval_onnx(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(s, WHITE)
    _add_header_bar(s, "Dense Retrieval & ONNX / INT8 Pipeline",
                    "FAISS for sub-millisecond lookup + dynamic quantisation for CPU edge deployment")
    # Left: FAISS card
    _add_rect(s, Inches(0.4), Inches(1.1), Inches(6.2), Inches(5.8), CREAM)
    _add_rect(s, Inches(0.4), Inches(1.1), Inches(6.2), Inches(0.55), DEEP_BLUE)
    _add_text(s, Inches(0.6), Inches(1.18), Inches(5.8), Inches(0.4),
              "FAISS IndexFlatIP", size=15, bold=True, color=WHITE, font=FONT_TITLE)
    _add_bullets(s, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.5), [
        ("All 8K gallery images pre-encoded into 256-d unit vectors (L2-normalised)", []),
        ("IndexFlatIP = Inner-Product on L2-normalised vectors ≡ cosine similarity", []),
        ("Exact nearest-neighbour (no approximation error) — fine for 8K corpus", []),
        ("FAISS lookup = 1 ms per query on a single CPU core", []),
        ("Thread-safe read-only accessor + @st.cache_resource → multiple concurrent Streamlit sessions don't race", []),
        ("Result: sub-second end-to-end retrieval, even on the 1 GB Streamlit Cloud sandbox", []),
    ], size=12, line_spacing=1.2)
    # Right: ONNX card
    _add_rect(s, Inches(6.7), Inches(1.1), Inches(6.2), Inches(5.8), WHITE, line=LIGHT_BLUE)
    _add_rect(s, Inches(6.7), Inches(1.1), Inches(6.2), Inches(0.55), DEEP_RED)
    _add_text(s, Inches(6.9), Inches(1.18), Inches(5.8), Inches(0.4),
              "ONNX + INT8 Quantisation", size=15, bold=True, color=WHITE, font=FONT_TITLE)
    _add_bullets(s, Inches(6.9), Inches(1.8), Inches(5.8), Inches(4.5), [
        ("Export PyTorch backbone → ONNX (opset 17, dynamic_axes for variable batch)", []),
        ("Why opset 17 not 18+: Streamlit Cloud ships fixed ONNX Runtime that doesn't support 18+", []),
        ("Dynamic INT8 quantisation: weights pre-quantised, activations quantised at runtime", []),
        ("Per-channel weight quantisation (scale + zero-point per output channel) → minimises L2 error", []),
        ("Deployment bundle = 15.5 MB (only the projection-head weights; backbone re-downloaded at launch)", []),
        ("Result: 3.1× speedup vs PyTorch FP32, <70 ms CPU end-to-end latency", []),
    ], size=12, line_spacing=1.2)
    _add_footer(s)
    _add_page_number(s, n, total)


def _results_patk(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(s, WHITE)
    _add_header_bar(s, "Results — Precision@K & Ablations",
                    "PEFT delivers +13.4 pp at P@1; every design choice in the ladder adds measurable gain")
    # Big numbers row
    nums = [
        ("+13.4 pp", "Precision@1 gain (fine-tuned 2-layer MLP vs pretrained CLIP baseline)", DEEP_RED),
        ("+18.8%",  "Relative gain at P@5 (84.6% vs 71.2%)", DEEP_BLUE),
        ("+10.5%",  "Relative gain at P@10 (91.4% vs 82.7%)", BRIGHT_RED),
        ("63.0%",   "Best P@1 (full pipeline: 2-layer MLP + GELU + dropout 0.1 + augmentation)", DEEP_BLUE),
    ]
    for i, (val, lbl, color) in enumerate(nums):
        x = Inches(0.4 + i * 3.2)
        _add_rect(s, x, Inches(1.1), Inches(3.0), Inches(1.6), color)
        _add_text(s, x, Inches(1.2), Inches(3.0), Inches(0.7), val,
                  size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                  anchor=MSO_ANCHOR.MIDDLE, font=FONT_TITLE)
        _add_text(s, x + Inches(0.15), Inches(1.95), Inches(2.7), Inches(0.7), lbl,
                  size=10, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_BODY)
    # Ablation table
    _add_text(s, Inches(0.4), Inches(3.0), Inches(12.5), Inches(0.4),
              "Ablation Ladder — Every Design Choice Adds Gain", size=14, bold=True, color=DEEP_BLUE, font=FONT_TITLE)
    abl = [
        ["Configuration", "Precision@1", "Δ vs. zero-shot"],
        ["(a) Zero-shot CLIP (no fine-tuning)",        "48.9%", "—"],
        ["(b) Fine-tuned 1-layer linear projection",   "58.6%", "+9.7 pp"],
        ["(c) Fine-tuned 2-layer MLP with GELU",        "61.1%", "+12.2 pp"],
        ["(d) Fine-tuned 2-layer MLP with ReLU",        "59.9%", "+11.0 pp"],
        ["(e) + Dropout 0.1",                           "62.3%", "+13.4 pp"],
        ["(f) + Data augmentation (H-flip + text rewrite)", "63.0%", "+14.1 pp"],
    ]
    table = s.shapes.add_table(len(abl), 3, Inches(0.4), Inches(3.45), Inches(12.5), Inches(2.8)).table
    for r, row in enumerate(abl):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = FONT_BODY
                    run.font.size = Pt(12)
                    run.font.bold = (r == 0) or (c > 0 and r == len(abl) - 1)
                    run.font.color.rgb = WHITE if r == 0 else (BRIGHT_RED if c > 0 and r == len(abl) - 1 else DARK_TEXT)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DEEP_BLUE
            elif r == len(abl) - 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CREAM
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CREAM
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    # Bottom takeaway
    _add_rect(s, Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.6), DEEP_BLUE)
    _add_text(s, Inches(0.6), Inches(6.45), Inches(12.1), Inches(0.5),
              "Takeaway: PEFT alone gives +9.7 pp (the single biggest jump) — every additional design choice compounds on top of it.",
              size=12, bold=True, color=WHITE, italic=True, font=FONT_BODY, anchor=MSO_ANCHOR.MIDDLE)
    _add_footer(s)
    _add_page_number(s, n, total)


def _results_latency(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(s, WHITE)
    _add_header_bar(s, "Results — Latency & Cross-Model Benchmark",
                    "Sub-100 ms CPU latency + ALIGN > BLIP-2 > CLIP on RF/telemetry fine-tuning")
    # Left: latency table
    _add_text(s, Inches(0.4), Inches(1.1), Inches(6.0), Inches(0.4),
              "End-to-End Query Latency (CPU, Batch=1, ViT-L/14)", size=14, bold=True, color=DEEP_BLUE, font=FONT_TITLE)
    lat = [
        ["Runtime", "Latency", "Speedup"],
        ["PyTorch FP32 (baseline)", "210 ms", "1.0×"],
        ["PyTorch FP16 (AMP)",      "~95 ms", "~2.2×"],
        ["ONNX + INT8 (deployed)",  "68 ms",  "3.1×"],
    ]
    table = s.shapes.add_table(len(lat), 3, Inches(0.4), Inches(1.6), Inches(6.0), Inches(2.0)).table
    for r, row in enumerate(lat):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = FONT_BODY
                    run.font.size = Pt(12)
                    run.font.bold = (r == 0) or (r == len(lat) - 1)
                    run.font.color.rgb = WHITE if r == 0 else (BRIGHT_RED if r == len(lat) - 1 else DARK_TEXT)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DEEP_BLUE
            elif r == len(lat) - 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CREAM
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    # Latency breakdown
    _add_text(s, Inches(0.4), Inches(3.9), Inches(6.0), Inches(0.4),
              "Per-Stage Breakdown (ONNX+INT8)", size=14, bold=True, color=DEEP_BLUE, font=FONT_TITLE)
    bd = [
        ("Tokenisation",       "2 ms"),
        ("Text encoder",       "28 ms"),
        ("Projection head",    "4 ms"),
        ("FAISS lookup",       "1 ms"),
        ("Image preprocessing","8 ms"),
        ("Image encoder",      "20 ms"),
        ("Streamlit UI render","5 ms"),
    ]
    for i, (k, v) in enumerate(bd):
        y = Inches(4.35 + i * 0.32)
        _add_rect(s, Inches(0.4), y, Inches(6.0), Inches(0.3), CREAM if i % 2 == 0 else WHITE)
        _add_text(s, Inches(0.55), y + Inches(0.02), Inches(4.5), Inches(0.28), k,
                  size=11, color=DARK_TEXT, font=FONT_BODY)
        _add_text(s, Inches(5.0), y + Inches(0.02), Inches(1.3), Inches(0.28), v,
                  size=11, bold=True, color=DEEP_BLUE, align=PP_ALIGN.RIGHT, font=FONT_TITLE)
    # Right: cross-model benchmark
    _add_text(s, Inches(6.9), Inches(1.1), Inches(6.0), Inches(0.4),
              "Cross-Model Benchmark (RF/telemetry val set, n=1600)", size=14, bold=True, color=DEEP_BLUE, font=FONT_TITLE)
    cm = [
        ["Model", "P@1 (Zero-shot)", "P@1 (Fine-tuned)"],
        ["CLIP  (ViT-B/32)",   "48.9%", "62.3%"],
        ["CLIP  (ViT-L/14)",   "52.4%", "66.1%"],
        ["ALIGN (ViT-H/14)",   "54.2%", "68.5%"],
        ["BLIP-2 (L/14+Q+OPT)","49.7%", "63.8%"],
    ]
    table = s.shapes.add_table(len(cm), 3, Inches(6.9), Inches(1.6), Inches(6.0), Inches(2.0)).table
    for r, row in enumerate(cm):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = FONT_BODY
                    run.font.size = Pt(11)
                    run.font.bold = (r == 0) or (r == 3)
                    run.font.color.rgb = WHITE if r == 0 else (BRIGHT_RED if r == 3 else DARK_TEXT)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DEEP_BLUE
            elif r == 3:  # ALIGN row highlighted
                cell.fill.solid()
                cell.fill.fore_color.rgb = CREAM
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    # Right takeaway
    _add_rect(s, Inches(6.9), Inches(3.9), Inches(6.0), Inches(3.0), DEEP_BLUE)
    _add_text(s, Inches(7.1), Inches(4.0), Inches(5.6), Inches(0.4),
              "Cross-Model Takeaways", size=14, bold=True, color=WHITE, font=FONT_TITLE)
    _add_bullets(s, Inches(7.1), Inches(4.5), Inches(5.6), Inches(2.4), [
        ("ALIGN (ViT-H/14) wins on P@1 — largest backbone, biggest pre-training corpus", []),
        ("BLIP-2 generates best captions but lower retrieval P@1 (Q-Former trained for generation, not discrimination)", []),
        ("CLIP ViT-B/32: best latency/quality trade-off → chosen for live demo", []),
    ], size=11, color=WHITE, bullet_color=BRIGHT_RED, line_spacing=1.2)
    _add_footer(s)
    _add_page_number(s, n, total)


def _results_dense_vs_sparse(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(s, WHITE)
    _add_header_bar(s, "Results — Dense vs Sparse Retrieval",
                    "On zero-keyword-overlap queries, dense latent decisively beats lexical baselines")
    # Example box
    _add_rect(s, Inches(0.4), Inches(1.1), Inches(12.5), Inches(1.3), CREAM, line=DEEP_RED)
    _add_text(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.4),
              "The Canonical Zero-Overlap Example", size=13, bold=True, color=DEEP_RED, font=FONT_TITLE)
    _add_text(s, Inches(0.6), Inches(1.6), Inches(12), Inches(0.4),
              "Query: \"a vivacious furry companion darting across a sandy waterfront\"",
              size=12, color=DARK_TEXT, italic=True, font=FONT_BODY)
    _add_text(s, Inches(0.6), Inches(1.95), Inches(12), Inches(0.4),
              "Gold caption: \"dog running on beach\"    →    0 keyword tokens in common",
              size=12, color=DARK_TEXT, italic=True, font=FONT_BODY)
    # Results table
    _add_text(s, Inches(0.4), Inches(2.6), Inches(12.5), Inches(0.4),
              "Precision@1 on Zero-Keyword-Overlap Queries", size=14, bold=True, color=DEEP_BLUE, font=FONT_TITLE)
    tbl = [
        ["Retrieval Method",                   "P@1",    "Vocab-Mismatch Robust"],
        ["CLIP  (Latent, Fine-Tuned)",         "62.3%",  "Yes"],
        ["BM25  (Lexical)",                    "18.2%",  "No"],
        ["TF-IDF (Lexical)",                   "12.5%",  "No"],
        ["CLIP  (strict-zero subset, n=120)",  "58.7%",  "Yes (truly zero overlap)"],
        ["BM25 / TF-IDF (strict-zero subset)", "0.0%",   "Complete failure"],
    ]
    table = s.shapes.add_table(len(tbl), 3, Inches(0.4), Inches(3.1), Inches(12.5), Inches(2.6)).table
    for r, row in enumerate(tbl):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = FONT_BODY
                    run.font.size = Pt(12)
                    run.font.bold = (r == 0) or (r == 1) or (r == 4)
                    run.font.color.rgb = WHITE if r == 0 else (BRIGHT_RED if r in (1, 4) else DARK_TEXT)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DEEP_BLUE
            elif r in (1, 4):
                cell.fill.solid()
                cell.fill.fore_color.rgb = CREAM
            elif r == 5:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GREY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    # Bottom takeaway
    _add_rect(s, Inches(0.4), Inches(5.9), Inches(12.5), Inches(1.1), DEEP_BLUE)
    _add_text(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.4),
              "Why This Matters", size=13, bold=True, color=WHITE, font=FONT_TITLE)
    _add_text(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.55),
              "Lexical retrieval (BM25, TF-IDF) is fundamentally a token-matching operation — it cannot bridge the vocabulary gap. Dense latent retrieval maps both query and image into a shared 256-d space where semantic neighbours are easily found. This is the entire motivation for VLM-based retrieval.",
              size=11, color=WHITE, italic=True, font=FONT_BODY)
    _add_footer(s)
    _add_page_number(s, n, total)


def _demo(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(s, WHITE)
    _add_header_bar(s, "Live Demo — 5 Tabs of the Streamlit App",
                    "Hosted at image-search-app-fxwnabswg8tdmbcy5syym3.streamlit.app")
    tabs = [
        ("Tab 1", "Text → Image",        DEEP_BLUE,  ["Search box + Top-K slider", "Optional ONNX mode (FP32/INT8)", "Cosine score bars per result"]),
        ("Tab 2", "Image → Caption",     DEEP_RED,   ["Upload any image", "BLIP-2 generates 1 natural-language caption", "Cross-check against Flickr8k captions"]),
        ("Tab 3", "Semantic vs Keyword", BRIGHT_RED, ["CLIP vs TF-IDF vs BM25 side-by-side", "Shows the dense-vs-sparse gap in action", "Same query → 3 different rankings"]),
        ("Tab 4", "CLIP vs BLIP vs ALIGN", LIGHT_BLUE,["3-column model comparison", "Per-model latency (ms)", "Per-model top-5 results"]),
        ("Tab 5", "Image → Image (Reverse Search)", DEEP_BLUE, ["Upload query → encode → find similar", "Optional text context (OOD fix)", "Low-score auto-warning for spectrograms"]),
    ]
    for i, (tab, title, color, items) in enumerate(tabs):
        x = Inches(0.4 + i * 2.55)
        _add_rect(s, x, Inches(1.1), Inches(2.4), Inches(0.5), color)
        _add_text(s, x, Inches(1.18), Inches(2.4), Inches(0.4),
                  tab, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_TITLE)
        _add_rect(s, x, Inches(1.6), Inches(2.4), Inches(2.0), CREAM)
        _add_text(s, x + Inches(0.1), Inches(1.7), Inches(2.2), Inches(0.4),
                  title, size=11, bold=True, color=DEEP_BLUE, font=FONT_TITLE)
        for j, it in enumerate(items):
            _add_text(s, x + Inches(0.1), Inches(2.1 + j * 0.4), Inches(2.2), Inches(0.4),
                      "• " + it, size=9, color=DARK_TEXT, font=FONT_BODY)
    # Datasets bar
    _add_rect(s, Inches(0.4), Inches(3.9), Inches(12.5), Inches(0.55), DEEP_BLUE)
    _add_text(s, Inches(0.6), Inches(3.97), Inches(12.1), Inches(0.4),
              "3 Datasets in the Demo",
              size=13, bold=True, color=WHITE, font=FONT_TITLE)
    ds_items = [
        ("Flickr8k (300-img subset)", "in-domain natural images", DEEP_RED),
        ("CC3M (1K pairs)",            "pre-training proxy",     DEEP_BLUE),
        ("OOD-RF (synthetic)",         "out-of-distribution test", BRIGHT_RED),
    ]
    for i, (name, desc, color) in enumerate(ds_items):
        x = Inches(0.4 + i * 4.17)
        _add_rect(s, x, Inches(4.55), Inches(4.0), Inches(0.9), color)
        _add_text(s, x + Inches(0.15), Inches(4.62), Inches(3.7), Inches(0.4),
                  name, size=12, bold=True, color=WHITE, font=FONT_TITLE)
        _add_text(s, x + Inches(0.15), Inches(4.97), Inches(3.7), Inches(0.4),
                  desc, size=10, color=WHITE, italic=True, font=FONT_BODY)
    # ONNX bar
    _add_rect(s, Inches(0.4), Inches(5.65), Inches(12.5), Inches(1.35), DEEP_BLUE)
    _add_text(s, Inches(0.6), Inches(5.7), Inches(12.1), Inches(0.4),
              "ONNX Mode (FP32 vs INT8)",
              size=13, bold=True, color=WHITE, font=FONT_TITLE)
    _add_bullets(s, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.9), [
        ("FP32 = full precision (slower, larger, baseline accuracy)", []),
        ("INT8 = dynamic quantised (3.1× faster, 4× smaller, accuracy preserved within 1-2%)", []),
    ], size=11, color=WHITE, bullet_color=BRIGHT_RED)
    _add_footer(s)
    _add_page_number(s, n, total)


def _wpr_bugs(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(s, WHITE)
    _add_header_bar(s, "WPR Bugs, Failures & Pending Work — Resolved",
                    "5 critical deploy failures hit between WPR9 and the viva; all 5 fixed and live on the cloud app")
    bugs = [
        ("Silent deploy crash on first user request",          "HF download triggered at first query → OOM in 1 GB sandbox",       "Added `download_models.py` with retry logic; pre-warm at startup"),
        ("packages.txt formatting rejected by Streamlit Cloud", "CRLF / extra blank lines in system dependencies",                  "Re-wrote with LF + validated via streamlit validation"),
        ("Python 3.14 (Cloud default) vs 3.11 (our venv)",     "open_clip not on 3.14 yet, onnxruntime 1.16 incompatibility",     "Added `.python-version` file → Cloud uses 3.11"),
        ("1 GB RAM OOM for BLIP/ALIGN on Cloud",               "Both models need ~5 GB together — way over the sandbox limit",    "Auto-skip via `os.path.isdir(\"/mount/src\")` detection; local still loads them"),
        ("Text token dtype cast to fp16 broke embedding",      "nn.Embedding requires Long indices, not Half",                     "Cast only image inputs; text tokens stay Long; added `_infer_visual_dtype()`"),
    ]
    _add_text(s, Inches(0.4), Inches(1.1), Inches(12.5), Inches(0.4),
              "5 Critical Fixes", size=14, bold=True, color=DEEP_BLUE, font=FONT_TITLE)
    tbl = [["#", "Bug", "Root Cause", "Fix"]] + [[f"{i+1}", b[0], b[1], b[2]] for i, b in enumerate(bugs)]
    table = s.shapes.add_table(len(tbl), 4, Inches(0.4), Inches(1.55), Inches(12.5), Inches(4.5)).table
    for r, row in enumerate(tbl):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c == 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = FONT_BODY
                    run.font.size = Pt(10)
                    run.font.bold = (r == 0)
                    run.font.color.rgb = WHITE if r == 0 else DARK_TEXT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DEEP_BLUE
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CREAM
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    # Set column widths
    table.columns[0].width = Inches(0.5)
    table.columns[1].width = Inches(3.2)
    table.columns[2].width = Inches(4.5)
    table.columns[3].width = Inches(4.3)
    # Bottom takeaway
    _add_rect(s, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.85), DEEP_RED)
    _add_text(s, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.4),
              "Bonus: 5 Pre-Generated Demo Spectrograms", size=12, bold=True, color=WHITE, font=FONT_TITLE)
    _add_text(s, Inches(0.6), Inches(6.65), Inches(12.1), Inches(0.35),
              "Top 80% original photo + bottom 20% viridis heatmap composite → Tab 5 retrieves original as top-1 with score ≈ 0.85. Clean demo without showing OOD limitation.",
              size=10, color=WHITE, italic=True, font=FONT_BODY)
    _add_footer(s)
    _add_page_number(s, n, total)


def _conclusion(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(s, WHITE)
    _add_header_bar(s, "Conclusion & Future Scope",
                    "6 directions, each grounded in the limitations identified in this work")
    # Left: Conclusion
    _add_rect(s, Inches(0.4), Inches(1.1), Inches(6.0), Inches(5.8), CREAM)
    _add_rect(s, Inches(0.4), Inches(1.1), Inches(6.0), Inches(0.5), DEEP_BLUE)
    _add_text(s, Inches(0.6), Inches(1.18), Inches(5.6), Inches(0.4),
              "Conclusion", size=14, bold=True, color=WHITE, font=FONT_TITLE)
    _add_bullets(s, Inches(0.6), Inches(1.75), Inches(5.6), Inches(5.0), [
        ("Combined 4 previously disjoint VLM subfields — PEFT, multi-modal prompt learning, post-training quantisation, dense retrieval — into one working system.", []),
        ("Demonstrated that the 4 fields are complementary, not antagonistic.", []),
        ("Achieved +13.4 pp Precision@1 over zero-shot baseline via projection-head-only fine-tuning (~0.1% params).", []),
        ("3.1× speedup and 15.5 MB deployment bundle via ONNX + INT8.", []),
        ("Sub-100 ms CPU end-to-end latency on commodity hardware.", []),
        ("Dense latent retrieval decisively beats BM25 / TF-IDF on zero-overlap queries (62.3% vs 18.2% / 12.5%).", []),
    ], size=12, line_spacing=1.15)
    # Right: Future scope
    _add_rect(s, Inches(6.7), Inches(1.1), Inches(6.2), Inches(5.8), WHITE, line=LIGHT_BLUE)
    _add_rect(s, Inches(6.7), Inches(1.1), Inches(6.2), Inches(0.5), BRIGHT_RED)
    _add_text(s, Inches(6.9), Inches(1.18), Inches(5.8), Inches(0.4),
              "6 Future Directions", size=14, bold=True, color=WHITE, font=FONT_TITLE)
    fs = [
        ("MobileCLIP distillation",   "→ edge devices (Apple Neural Engine, RPi 5)"),
        ("SigLIP loss + bigger batch","→ +2-4 pp Precision@1 at batch=2K+"),
        ("Video & 3D modalities",     "→ temporal aggregation on frame embeddings"),
        ("Continual learning",        "→ replay buffer for streaming telemetry"),
        ("Hybrid sparse+dense",       "→ BM25 first-stage + CLIP re-ranker"),
        ("Active learning",            "→ cold-start domains (satellite, medical)"),
    ]
    for i, (k, v) in enumerate(fs):
        y = Inches(1.75 + i * 0.85)
        _add_rect(s, Inches(6.7), y, Inches(6.2), Inches(0.75), CREAM if i % 2 == 0 else WHITE)
        _add_rect(s, Inches(6.7), y, Inches(0.12), Inches(0.75), BRIGHT_RED)
        _add_text(s, Inches(6.95), y + Inches(0.1), Inches(5.9), Inches(0.3),
                  f"{i+1}.  {k}", size=12, bold=True, color=DEEP_BLUE, font=FONT_TITLE)
        _add_text(s, Inches(6.95), y + Inches(0.4), Inches(5.9), Inches(0.3),
                  v, size=10, color=GREY_TEXT, italic=True, font=FONT_BODY)
    # Bottom: target
    _add_rect(s, Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.45), DEEP_BLUE)
    _add_text(s, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.35),
              "Next-iteration targets: <30 ms latency  •  >75% Precision@1  •  >100K image-text pairs across domains",
              size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_TITLE)
    _add_footer(s)
    _add_page_number(s, n, total)


def _thanks(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(s, DEEP_BLUE)
    # Top accent
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.5), DEEP_RED)
    # Big title
    _add_text(s, Inches(0.7), Inches(1.5), Inches(12), Inches(1.2),
              "Thank You", size=72, bold=True, color=WHITE,
              align=PP_ALIGN.CENTER, font=FONT_TITLE)
    _add_text(s, Inches(0.7), Inches(2.7), Inches(12), Inches(0.6),
              "Questions & Discussion", size=24, color=LIGHT_BLUE,
              align=PP_ALIGN.CENTER, italic=True, font=FONT_TITLE)
    # Divider
    _add_rect(s, Inches(6.0), Inches(3.6), Inches(1.333), Inches(0.05), BRIGHT_RED)
    # Code & Data
    _add_text(s, Inches(0.7), Inches(4.0), Inches(12), Inches(0.4),
              "Code & Live App", size=14, bold=True, color=LIGHT_BLUE,
              align=PP_ALIGN.CENTER, font=FONT_TITLE)
    _add_text(s, Inches(0.7), Inches(4.4), Inches(12), Inches(0.4),
              "github.com/pd877083/image-search-app",
              size=16, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_MONO)
    _add_text(s, Inches(0.7), Inches(4.75), Inches(12), Inches(0.4),
              "image-search-app-fxwnabswg8tdmbcy5syym3.streamlit.app",
              size=16, color=WHITE, align=PP_ALIGN.CENTER, font=FONT_MONO)
    # References list
    _add_text(s, Inches(0.7), Inches(5.5), Inches(12), Inches(0.4),
              "Key References", size=14, bold=True, color=LIGHT_BLUE,
              align=PP_ALIGN.CENTER, font=FONT_TITLE)
    refs = "CLIP-Adapter (Gao 2021)  •  BLIP-2 (Li 2023)  •  CoOp (Zhou 2022)  •  AudioCLIP (Guzhov 2022)  •  MobileCLIP (Vasu 2024)"
    _add_text(s, Inches(0.7), Inches(5.9), Inches(12), Inches(0.4), refs,
              size=11, color=LIGHT_BLUE, align=PP_ALIGN.CENTER, font=FONT_BODY)
    refs2 = "SigLIP (Zhai 2023)  •  MaPLe (Khattak 2023)  •  Q-VLM (Wang 2024)  •  MARVEL (Zhou 2024)  •  Don't Stop Learning (Ding 2023)"
    _add_text(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.4), refs2,
              size=11, color=LIGHT_BLUE, align=PP_ALIGN.CENTER, font=FONT_BODY)
    # Author
    _add_text(s, Inches(0.7), Inches(6.9), Inches(12), Inches(0.4),
              "Dhruv Singhal  •  A2305224486  •  Under Dr. Rakesh Chandra Joshi  •  Amity University, Noida",
              size=12, color=BRIGHT_RED, italic=True, align=PP_ALIGN.CENTER, font=FONT_BODY)
    _add_page_number(s, n, total)


# ── Main ───────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = [
        _cover,
        _outline,
        _motivation,
        _objectives,
        _related_work,
        _datasets,
        _methodology,
        _peft,
        _retrieval_onnx,
        _results_patk,
        _results_latency,
        _results_dense_vs_sparse,
        _demo,
        _wpr_bugs,
        _conclusion,
        _thanks,
    ]
    total = len(slides)
    for i, build in enumerate(slides, start=1):
        if i == 1:
            build(prs)  # cover
        else:
            build(prs, i, total)

    prs.save(str(OUT_PATH))
    print(f"PPTX written: {OUT_PATH}")
    print(f"Size: {OUT_PATH.stat().st_size / 1024:.1f} KB")
    print(f"Slides: {total}")


if __name__ == "__main__":
    main()
